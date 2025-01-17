import os
import sys
from datetime import datetime
import pptx
import pdfplumber
from docx import Document
import logging
from colorama import Fore, Style, init
import re
import argparse
from pptx import Presentation
import zipfile
import tempfile
import shutil

# Initialize colorama
init(autoreset=True)
# Set up logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

def sanitize_text(text):
    """Removes non-printable characters and ensures proper markdown escaping."""
    # First remove non-printable characters
    text = re.sub(r'[\x00-\x1F\x7F-\x9F]', ' ', text)
    # Escape markdown special characters
    markdown_chars = ['#', '*', '_', '`', '[', ']', '(', ')', '<', '>', '|']
    for char in markdown_chars:
        text = text.replace(char, '\\' + char)
    return text

def setup_file_handlers():
    """Configure file type handlers and their processing methods."""
    
    # Define supported file types and their processors
    extract_functions = {
        '.pdf': extract_pdf_content,
        '.pptx': extract_pptx_content,
        '.docx': extract_docx_content,
        '.txt': extract_text_content,
        '.py': extract_code_content,
        '.java': extract_code_content,
        '.js': extract_code_content,
        '.cpp': extract_code_content,
        '.sql': extract_code_content,
        '.json': extract_code_content,
        '.yaml': extract_code_content,
        '.yml': extract_code_content,
        '.xml': extract_code_content,
        '.html': extract_code_content,
        '.css': extract_code_content,
        '.md': extract_code_content
    }
    
    # Define language identifiers for code highlighting
    code_file_types = {
        '.py': 'python',
        '.java': 'java',
        '.js': 'javascript',
        '.cpp': 'cpp',
        '.sql': 'sql',
        '.json': 'json',
        '.yaml': 'yaml',
        '.yml': 'yaml',
        '.xml': 'xml',
        '.html': 'html',
        '.css': 'css',
        '.md': 'markdown'
    }
    
    return extract_functions, code_file_types

def extract_pdf_content(filepath):
    """Extract text from PDF files."""
    try:
        with pdfplumber.open(filepath) as pdf:
            return "\n".join(page.extract_text() for page in pdf.pages if page.extract_text())
    except Exception as e:
        logging.warning(f"Error extracting PDF content from {filepath}: {str(e)}")
        return f"Error extracting PDF content: {str(e)}"

def extract_pptx_content(filepath):
    """Extract text from PowerPoint presentations."""
    try:
        prs = Presentation(filepath)
        text_content = []
        
        for i, slide in enumerate(prs.slides, 1):
            text_content.append(f"\n## Slide {i}")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)
        
        return "\n".join(text_content)
    except Exception as e:
        logging.warning(f"Error extracting PPTX content from {filepath}: {str(e)}")
        return f"Error extracting PPTX content: {str(e)}"

def extract_docx_content(filepath):
    """Extract text from Word documents."""
    try:
        doc = Document(filepath)
        return "\n".join(paragraph.text for paragraph in doc.paragraphs)
    except Exception as e:
        logging.warning(f"Error extracting DOCX content from {filepath}: {str(e)}")
        return f"Error extracting DOCX content: {str(e)}"

def extract_text_content(filepath):
    """Extract content from text files."""
    try:
        with open(filepath, 'r', encoding='utf-8') as file:
            return file.read()
    except UnicodeDecodeError:
        try:
            with open(filepath, 'r', encoding='latin-1') as file:
                return file.read()
        except Exception as e:
            logging.warning(f"Error reading text file {filepath}: {str(e)}")
            return f"Error reading text file: {str(e)}"

def extract_code_content(filepath):
    """Extract content from source code files."""
    try:
        with open(filepath, 'r', encoding='utf-8') as file:
            return file.read()
    except Exception as e:
        logging.warning(f"Error reading code file {filepath}: {str(e)}")
        return f"Error reading code file: {str(e)}"

def detect_file_type(filepath, content):
    """Enhanced file type detection."""
    filename = os.path.basename(filepath).lower()
    content_sample = content[:1000].lower() if content else ""
    
    # Medical reports
    if "zpráva" in filename or any(marker in content_sample for marker in ["LÉKAŘSKÁ ZPRÁVA", "Pacient:", "Rodné číslo:"]):
        return "medical_report"
    
    # Configuration files
    elif filename.endswith(('.json', '.yaml', '.yml', '.toml', '.ini', '.conf')):
        return "config"
    
    # Source code files
    elif filename.endswith(('.py', '.js', '.java', '.cpp', '.cs', '.php')):
        return "source_code"
    
    # SQL files or content
    elif filename.endswith('.sql') or any(marker in content_sample for marker in ["SELECT", "INSERT", "UPDATE", "CREATE TABLE"]):
        return "sql"
    
    # Log files
    elif filename.endswith('.log') or any(marker in content_sample for marker in ["ERROR", "WARNING", "INFO", "[DEBUG]"]):
        return "log"
    
    # Markdown/Documentation
    elif filename.endswith(('.md', '.rst', '.txt')):
        return "documentation"
    
    return "default"

def format_medical_report(content):
    """Formats medical reports in a structured way."""
    formatted_parts = []
    lines = content.strip().split('\n')
    
    formatted_parts.append("# LÉKAŘSKÁ ZPRÁVA\n")
    
    # Patient Information Table
    formatted_parts.append("## Informace o pacientovi")
    formatted_parts.append("| Položka | Hodnota |")
    formatted_parts.append("|---------|----------|")
    
    patient_fields = ["Pacient:", "Rodné číslo:", "Bydliště:", "Pojišťovna:", 
                     "Telefon:", "Věk:", "Datum:"]
    
    for line in lines:
        for field in patient_fields:
            if field in line:
                value = line.split(field)[1].strip()
                formatted_parts.append(f"| {field[:-1]} | {value} |")
    
    # Results Section
    formatted_parts.append("\n## Výsledky vyšetření")
    formatted_parts.append("```")
    in_results = False
    for line in lines:
        if "Laboratorní výsledky" in line or "Dg.:" in line:
            in_results = True
        if in_results and line.strip():
            formatted_parts.append(line.strip())
    formatted_parts.append("```")
    
    return "\n".join(formatted_parts)

def process_file_content(filepath, content, ext, code_file_types):
    """Process file content with appropriate formatting."""
    formatted_content = []
    
    # Add file metadata
    file_size = os.path.getsize(filepath)
    last_modified = datetime.fromtimestamp(os.path.getmtime(filepath))
    
    formatted_content.append("### Metadata")
    formatted_content.append("| Attribute | Value |")
    formatted_content.append("|-----------|--------|")
    formatted_content.append(f"| File Size | {file_size:,} bytes |")
    formatted_content.append(f"| Last Modified | {last_modified.strftime('%Y-%m-%d %H:%M:%S')} |")
    formatted_content.append(f"| File Type | {ext} |")
    formatted_content.append("\n### Content\n")
    
    # Detect and format content based on type
    content_type = detect_file_type(filepath, content)
    
    if content_type == "medical_report":
        formatted_content.append(format_medical_report(content))
    elif ext in code_file_types:
        # Code files get syntax highlighting
        formatted_content.append(f"```{code_file_types[ext]}")
        formatted_content.append(content)
        formatted_content.append("```")
    else:
        # Default text formatting with escaping
        formatted_content.append("```")
        formatted_content.append(content)
        formatted_content.append("```")
    
    return "\n".join(formatted_content) + "\n\n"

def build_folder_structure_markdown(root_path, excluded_folders):
    """Builds a tree-like markdown structure of all files & folders."""
    structure_lines = []
    
    for current_root, dirs, files in os.walk(root_path):
        # Remove excluded directories
        dirs[:] = [d for d in dirs if d not in excluded_folders]
        
        level = current_root.replace(root_path, '').count(os.sep)
        indent = "  " * level
        folder_name = os.path.basename(current_root)
        if not folder_name:
            folder_name = os.path.basename(root_path)

        structure_lines.append(f"{indent}- **{folder_name}/**")

        for f in files:
            file_indent = "  " * (level + 1)
            rel_path = os.path.relpath(os.path.join(current_root, f), root_path)
            structure_lines.append(f"{file_indent}- [{f}](#{rel_path.replace(os.sep,'-')})")

    return "\n".join(structure_lines)

def process_zip_file(zip_path):
    """Extract and process a zip file."""
    temp_dir = tempfile.mkdtemp()
    try:
        with zipfile.ZipFile(zip_path, 'r') as zip_ref:
            zip_ref.extractall(temp_dir)
        
        # Process the extracted contents using existing functionality
        folder_name = os.path.basename(zip_path).replace('.zip', '')
        
        # Process the temp directory using process_folder instead of process_folder_contents
        output_md = process_folder(temp_dir)
        
        return output_md
    finally:
        # Clean up temp directory
        shutil.rmtree(temp_dir)

def process_folder(folder_path, excluded_folders=None):
    """Process a folder and create a markdown file with its contents."""
    if excluded_folders is None:
        excluded_folders = ["venv", "__pycache__", ".git"]
    
    # Setup file handlers
    extract_functions, code_file_types = setup_file_handlers()
    
    # Generate output filename
    folder_name = os.path.basename(folder_path)
    output_md = f"combined_{folder_name}.md"
    
    # Process the folder contents
    with open(output_md, 'w', encoding='utf-8') as outfile:
        # Write header
        outfile.write(f"# Contents of {folder_name}\n\n")
        
        # Add folder structure
        outfile.write("## Folder Structure\n")
        folder_structure = build_folder_structure_markdown(folder_path, excluded_folders)
        outfile.write(folder_structure + "\n\n")
        
        # Process each file
        outfile.write("## File Contents\n\n")
        for root, dirs, files in os.walk(folder_path):
            # Skip excluded directories
            dirs[:] = [d for d in dirs if d not in excluded_folders]
            
            for file in files:
                filepath = os.path.join(root, file)
                ext = os.path.splitext(file)[1].lower()
                
                if ext in extract_functions:
                    try:
                        # Extract content using appropriate function
                        content = extract_functions[ext](filepath)
                        if content:
                            # Create anchor for file
                            rel_path = os.path.relpath(filepath, folder_path)
                            outfile.write(f"<a id='{rel_path.replace(os.sep, '-')}'></a>\n")
                            outfile.write(f"### {rel_path}\n\n")
                            
                            # Process and write content
                            formatted_content = process_file_content(filepath, content, ext, code_file_types)
                            outfile.write(formatted_content)
                    except Exception as e:
                        logging.error(f"Error processing file {filepath}: {str(e)}")
    
    return output_md

def main():
    parser = argparse.ArgumentParser(description="Combine folder or zip contents into a single Markdown file.")
    parser.add_argument("path", help="Path to the folder or zip file to process.")
    parser.add_argument("--exclude", nargs="+", default=["venv", "__pycache__", ".git"],
                      help="Folders to exclude from processing (default: venv __pycache__ .git)")
    args = parser.parse_args()
    
    input_path = os.path.abspath(args.path)
    
    if zipfile.is_zipfile(input_path):
        logging.info(f"Processing zip file: {input_path}")
        output_file = process_zip_file(input_path)
        logging.info(f"Zip file processed. Output saved to '{output_file}'")
    else:
        # Existing folder processing logic
        folder_name = os.path.basename(input_path)
        output_md = f"combined_{folder_name}.md"
        extract_functions, code_file_types = setup_file_handlers()
        process_folder_contents(input_path, output_md, extract_functions, code_file_types)
        logging.info(f"Folder processed. Output saved to '{output_md}'")

if __name__ == "__main__":
    main() 